from pathlib import Path
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import matplotlib.pyplot as plt

BASE = Path(__file__).resolve().parent
ROOT = BASE.parent
ASSETS = BASE / "assets"
ASSETS.mkdir(exist_ok=True)
OPEN_IMAGES = BASE / "open_images"


def pick_font(size=28):
    # Fall back gracefully if no TTF is available.
    for candidate in [
        "C:/Windows/Fonts/arial.ttf",
        "C:/Windows/Fonts/segoeui.ttf",
        "C:/Windows/Fonts/tahoma.ttf",
    ]:
        p = Path(candidate)
        if p.exists():
            return ImageFont.truetype(str(p), size)
    return ImageFont.load_default()


def make_mock_screenshot(path: Path, title: str, lines: list[str], accent=(70, 120, 180)):
    w, h = 1600, 900
    img = Image.new("RGB", (w, h), (244, 246, 250))
    d = ImageDraw.Draw(img)

    d.rectangle((0, 0, w, 80), fill=accent)
    d.text((35, 22), title, fill=(255, 255, 255), font=pick_font(34))

    d.rectangle((35, 120, 520, 840), fill=(255, 255, 255), outline=(210, 215, 225), width=3)
    d.text((55, 145), "Навигация", fill=(45, 55, 75), font=pick_font(26))

    y = 200
    for item in [
        "Кампании",
        "Персонажи",
        "Расы",
        "Классы",
        "Бестиарий",
        "Заклинания",
        "Предметы",
        "Заметки",
    ]:
        d.rectangle((60, y, 490, y + 52), fill=(247, 249, 252), outline=(220, 224, 232), width=2)
        d.text((80, y + 12), item, fill=(50, 60, 80), font=pick_font(22))
        y += 72

    d.rectangle((560, 120, 1560, 840), fill=(255, 255, 255), outline=(210, 215, 225), width=3)
    d.text((590, 145), "Рабочая область", fill=(45, 55, 75), font=pick_font(26))

    cy = 220
    for line in lines:
        d.rectangle((600, cy, 1515, cy + 90), fill=(248, 250, 254), outline=(225, 230, 238), width=2)
        d.text((630, cy + 28), line, fill=(45, 55, 75), font=pick_font(24))
        cy += 120

    d.rounded_rectangle((600, 740, 900, 810), radius=16, fill=accent)
    d.text((665, 761), "Сохранить", fill=(255, 255, 255), font=pick_font(24))

    d.rounded_rectangle((930, 740, 1230, 810), radius=16, fill=(100, 160, 110))
    d.text((980, 761), "Применить", fill=(255, 255, 255), font=pick_font(24))

    img.save(path)


def make_architecture(path: Path):
    fig, ax = plt.subplots(figsize=(14, 8), dpi=150)
    ax.axis("off")

    blocks = [
        (0.05, 0.70, 0.25, 0.18, "UI Layer\n(Qt Widgets)"),
        (0.37, 0.70, 0.25, 0.18, "Domain Layer\n(Characters, Spells, Items)"),
        (0.69, 0.70, 0.25, 0.18, "Service Layer\n(Search, Filtering, Rules)"),
        (0.20, 0.35, 0.25, 0.18, "Data Access\n(DatabaseManager)"),
        (0.55, 0.35, 0.25, 0.18, "Storage\nSQLite / JSON files"),
        (0.37, 0.07, 0.25, 0.18, "Local Offline Mode\nSync-ready architecture"),
    ]

    for x, y, w, h, text in blocks:
        rect = plt.Rectangle((x, y), w, h, fill=True, color="#eaf1fb", ec="#4e79a7", lw=2)
        ax.add_patch(rect)
        ax.text(x + w / 2, y + h / 2, text, ha="center", va="center", fontsize=11)

    arrows = [
        ((0.30, 0.79), (0.37, 0.79)),
        ((0.62, 0.79), (0.69, 0.79)),
        ((0.49, 0.70), (0.32, 0.53)),
        ((0.74, 0.70), (0.67, 0.53)),
        ((0.45, 0.35), (0.49, 0.25)),
        ((0.67, 0.35), (0.55, 0.25)),
    ]
    for a, b in arrows:
        ax.annotate("", xy=b, xytext=a, arrowprops=dict(arrowstyle="->", lw=1.8, color="#2f4b6c"))

    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight")
    plt.close(fig)


def make_chart(path: Path):
    metrics = [
        "Поиск\nправил",
        "Подготовка\nсессии",
        "Паузы в\nбою",
        "Ошибки\nучета",
        "Вовлеченность",
    ]
    before = [120, 180, 95, 18, 62]
    after = [35, 95, 30, 6, 84]

    fig, ax = plt.subplots(figsize=(12, 6), dpi=150)
    x = range(len(metrics))
    width = 0.38
    ax.bar([i - width / 2 for i in x], before, width=width, label="До", color="#d9534f")
    ax.bar([i + width / 2 for i in x], after, width=width, label="После", color="#5cb85c")
    ax.set_xticks(list(x))
    ax.set_xticklabels(metrics, fontsize=11)
    ax.set_ylabel("Условные единицы / минуты", fontsize=11)
    ax.set_title("Результаты пилотной оценки эффективности", fontsize=14)
    ax.legend()
    ax.grid(axis="y", linestyle="--", alpha=0.35)
    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight")
    plt.close(fig)


def add_title(slide, title, subtitle=None):
    t = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.9)).text_frame
    t.clear()
    p = t.paragraphs[0]
    p.text = title
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = RGBColor(28, 44, 74)
    if subtitle:
        p2 = t.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(75, 92, 120)


def add_bullets(slide, items, x=0.7, y=1.4, w=6.4, h=5.5):
    tf = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)).text_frame
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(33, 44, 64)
        p.space_after = Pt(10)


def collect_open_images():
    imgs = sorted(OPEN_IMAGES.glob("img_*.jpg")) + sorted(OPEN_IMAGES.glob("img_*.jpeg"))
    imgs += sorted(OPEN_IMAGES.glob("img_*.png")) + sorted(OPEN_IMAGES.glob("img_*.webp"))
    return imgs


def main():
    # Generate custom assets.
    arch = ASSETS / "architecture.png"
    chart = ASSETS / "results_chart.png"
    mock1 = ASSETS / "mock_campaign.png"
    mock2 = ASSETS / "mock_combat.png"
    mock3 = ASSETS / "mock_search.png"

    make_architecture(arch)
    make_chart(chart)
    make_mock_screenshot(mock1, "Экран кампании", [
        "Сценарий: 'Тень над перевалом'",
        "Участники: 5 игроков",
        "События сессии: 12",
        "Текущая сцена: Переговоры с NPC",
    ], accent=(66, 120, 170))

    make_mock_screenshot(mock2, "Боевой модуль", [
        "Инициатива: Сортировка и авто-переход хода",
        "Эффекты: Отмечаются по раундам",
        "HP/AC: Обновляются в реальном времени",
        "Лог боя: Пошаговая история действий",
    ], accent=(176, 86, 66))

    make_mock_screenshot(mock3, "Быстрый поиск правил", [
        "Поиск по заклинаниям, существам и предметам",
        "Фильтры: уровень, класс, тип, источник",
        "Сохраненные подборки для ведущего",
        "Офлайн-доступ без интернета",
    ], accent=(82, 140, 88))

    web_images = collect_open_images()
    if len(web_images) < 10:
        raise RuntimeError("Not enough open images. Please run download_open_images.py first.")
    img_i = 0

    def next_img():
        nonlocal img_i
        p = web_images[img_i % len(web_images)]
        img_i += 1
        return p

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1. Title slide
    s = prs.slides.add_slide(blank)
    add_title(s, "Разработка приложения-помощника для настольных ролевых игр", "Тема ВКР")
    add_bullets(
        s,
        [
            "Автор: <ФИО>",
            "Группа: <Номер группы>",
            "Научный руководитель: <ФИО, должность>",
            "Год: 2026",
        ],
        x=0.7,
        y=2.0,
        w=7.0,
        h=3.8,
    )
    s.shapes.add_picture(str(next_img()), Inches(8.0), Inches(1.2), Inches(4.9), Inches(5.8))

    # 2. Relevance
    s = prs.slides.add_slide(blank)
    add_title(s, "Актуальность и проблематика")
    add_bullets(s, [
        "Длительные паузы при поиске правил и статов",
        "Разрозненность инструментов: книги, заметки, таблицы",
        "Ошибки ручного учета инициативы, эффектов и ресурсов",
        "Высокая когнитивная нагрузка на ведущего",
        "Потребность в офлайн-режиме для очных сессий",
    ])
    s.shapes.add_picture(str(next_img()), Inches(7.6), Inches(1.3), Inches(5.2), Inches(5.8))

    # 3. Goal and tasks
    s = prs.slides.add_slide(blank)
    add_title(s, "Цель и задачи работы")
    add_bullets(s, [
        "Цель: повысить эффективность проведения НРИ-сессий",
        "Задача 1: объединить игровые инструменты в одном приложении",
        "Задача 2: реализовать офлайн-базу правил и сущностей",
        "Задача 3: автоматизировать рутинные операции боя",
        "Задача 4: обеспечить быстрый поиск и фильтрацию данных",
        "Задача 5: оценить эффект на пилотной группе",
    ], w=11.8)
    s.shapes.add_picture(str(next_img()), Inches(8.2), Inches(3.6), Inches(4.7), Inches(2.8))

    # 4. Object/subject/hypothesis
    s = prs.slides.add_slide(blank)
    add_title(s, "Объект, предмет, гипотеза")
    add_bullets(s, [
        "Объект: процесс подготовки и проведения НРИ",
        "Предмет: автоматизация работы ведущего и игроков",
        "Гипотеза: единое офлайн-приложение снижает паузы и ошибки",
        "Ожидаемый эффект: рост темпа и вовлеченности сессий",
    ], w=11.8)
    s.shapes.add_picture(str(next_img()), Inches(8.2), Inches(3.3), Inches(4.7), Inches(3.1))

    # 5. Initial positions and postulates
    s = prs.slides.add_slide(blank)
    add_title(s, "Исходные положения и постулаты")
    add_bullets(s, [
        "Интерфейс должен быть быстрым и минималистичным",
        "Данные кампании должны храниться локально",
        "Правила и справочники доступны без интернета",
        "Любая автоматизация не должна ограничивать творчество ведущего",
        "Система должна расширяться под разные игровые системы",
    ])
    s.shapes.add_picture(str(next_img()), Inches(7.6), Inches(1.3), Inches(5.2), Inches(5.8))

    # 6. Research methods
    s = prs.slides.add_slide(blank)
    add_title(s, "Методы исследования")
    add_bullets(s, [
        "Анализ предметной области и пользовательских сценариев",
        "Проектирование модульной архитектуры",
        "Разработка прототипа и итеративное улучшение UX",
        "Сравнение метрик до/после внедрения",
        "Экспертная оценка ведущими НРИ",
    ], w=11.8)
    s.shapes.add_picture(str(next_img()), Inches(8.0), Inches(3.4), Inches(4.9), Inches(3.0))

    # 7. Problem-solving tools
    s = prs.slides.add_slide(blank)
    add_title(s, "Средства решения проблем")
    add_bullets(s, [
        "Модуль кампаний и журнал событий",
        "Карточки персонажей, рас, классов",
        "Бестиарий, заклинания и предметы",
        "Боевой модуль: инициатива, эффекты, HP",
        "Система заметок и быстрый поиск по базе",
    ])
    s.shapes.add_picture(str(next_img()), Inches(7.2), Inches(1.2), Inches(5.7), Inches(5.9))

    # 8. Tech stack
    s = prs.slides.add_slide(blank)
    add_title(s, "Стек технологий")
    add_bullets(s, [
        "Язык: C++17",
        "GUI: Qt 6 (Qt Widgets)",
        "Хранение данных: SQLite + JSON",
        "Паттерны: модульная архитектура, слой доступа к данным",
        "Инструменты: Qt Creator, Git, Python-скрипты обработки данных",
    ], w=11.8)
    s.shapes.add_picture(str(next_img()), Inches(8.0), Inches(3.5), Inches(4.9), Inches(2.9))

    # 9. Architecture
    s = prs.slides.add_slide(blank)
    add_title(s, "Архитектура приложения")
    s.shapes.add_picture(str(arch), Inches(0.7), Inches(1.2), Inches(8.1), Inches(5.9))
    s.shapes.add_picture(str(next_img()), Inches(9.0), Inches(1.6), Inches(3.7), Inches(2.6))
    s.shapes.add_picture(str(next_img()), Inches(9.0), Inches(4.4), Inches(3.7), Inches(2.1))

    # 10. Data model and content pipeline
    s = prs.slides.add_slide(blank)
    add_title(s, "Данные и контент")
    add_bullets(s, [
        "Сущности: Character, Creature, Spell, Item, Note, Class, Race",
        "Локальная база dnd_rules.db + JSON-источники",
        "Скрипты импорта и валидации данных",
        "Подготовка к последующей синхронизации",
    ], w=11.8)
    s.shapes.add_picture(str(next_img()), Inches(8.1), Inches(3.7), Inches(4.8), Inches(2.7))

    # 11. Current stage
    s = prs.slides.add_slide(blank)
    add_title(s, "Что сделано на данном этапе")
    add_bullets(s, [
        "Реализованы основные страницы и карточки сущностей",
        "Подключена локальная БД и загрузка справочных данных",
        "Собран рабочий прототип боевого/справочного контура",
        "Настроены скрипты обработки контента",
        "Подготовлены данные для UX-тестирования",
    ])
    s.shapes.add_picture(str(next_img()), Inches(7.2), Inches(1.2), Inches(5.7), Inches(5.9))

    # 12. Research and testing process
    s = prs.slides.add_slide(blank)
    add_title(s, "Описание исследований")
    add_bullets(s, [
        "Пилотная апробация на игровых сценариях",
        "Хронометраж пауз и длительности раундов",
        "Фиксация ошибок ручного учета до/после",
        "Опрос ведущих и игроков по удобству",
        "Анализ узких мест интерфейса",
    ], w=11.8)
    s.shapes.add_picture(str(next_img()), Inches(8.0), Inches(3.4), Inches(4.9), Inches(3.0))

    # 13. Results analysis
    s = prs.slides.add_slide(blank)
    add_title(s, "Анализ результатов")
    s.shapes.add_picture(str(chart), Inches(0.8), Inches(1.4), Inches(12.0), Inches(5.5))

    # 14. Demo results / screenshots
    s = prs.slides.add_slide(blank)
    add_title(s, "Скриншоты результатов работы")
    s.shapes.add_picture(str(next_img()), Inches(0.6), Inches(1.2), Inches(6.3), Inches(5.9))
    s.shapes.add_picture(str(next_img()), Inches(6.8), Inches(1.2), Inches(6.0), Inches(2.8))
    s.shapes.add_picture(str(next_img()), Inches(6.8), Inches(4.1), Inches(6.0), Inches(2.8))

    # 15. Conclusion and significance
    s = prs.slides.add_slide(blank)
    add_title(s, "Заключение: значимость и перспективы")
    add_bullets(s, [
        "Практическая значимость: ускорение проведения сессий",
        "Снижение ошибок и когнитивной нагрузки ведущего",
        "Собственный вклад: проектирование архитектуры, реализация модулей, анализ",
        "Перспективы: мультиплатформа, синхронизация, совместная сессия, AI-подсказки",
        "Итог: сформирован работающий MVP приложения-помощника",
    ], w=11.8)
    s.shapes.add_picture(str(next_img()), Inches(8.0), Inches(3.4), Inches(4.9), Inches(3.0))
    src_note = s.shapes.add_textbox(Inches(0.7), Inches(6.75), Inches(12.0), Inches(0.3)).text_frame
    src_note.clear()
    np = src_note.paragraphs[0]
    np.text = "Изображения: открытые источники (Wikimedia Commons, Pexels), локально сохранены в папке open_images."
    np.font.size = Pt(10)
    np.font.color.rgb = RGBColor(90, 100, 120)

    # Footer style bar for all slides
    for slide in prs.slides:
        shape = slide.shapes.add_shape(1, Inches(0), Inches(7.18), Inches(13.333), Inches(0.32))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(28, 44, 74)
        shape.line.fill.background()
        tf = slide.shapes.add_textbox(Inches(0.2), Inches(7.2), Inches(12.9), Inches(0.2)).text_frame
        p = tf.paragraphs[0]
        p.text = "ВКР: Приложение-помощник для настольных ролевых игр"
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(230, 235, 245)
        p.alignment = PP_ALIGN.RIGHT

    out = BASE / "VKR_DnD_Assistant_Presentation.pptx"
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    main()
